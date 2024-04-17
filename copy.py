import pptx
from pptx import Presentation

ppt = Presentation("Test.pptx")

import copy

def copy_slide_from_external_prs(prs, location, slide_no):
    """
    prs = ppt object to which slides has to be copied
    location = location of the ppt from which the slides has to be copied
    slide_no = no. of the slide that has to be copied
    """
    # copy from external presentation all objects into the existing presentation
    external_pres = Presentation(location)

    # specify the slide you want to copy the contents from
    ext_slide = external_pres.slides[slide_no]

    # Define the layout you want to use from your generated pptx
    slide_layout = prs.slide_layouts[slide_no]

    # create now slide, to copy contents to 
    curr_slide = prs.slides.add_slide(slide_layout)

    # now copy contents from external slide, but do not copy slide properties
    # e.g. slide layouts, etc., because these would produce errors, as diplicate
    # entries might be generated

    for shp in ext_slide.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        curr_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    return prs

ppt = copy_slide_from_external_prs(ppt,"Test_modified.pptx", 0)

ppt.save("mod.pptx")
